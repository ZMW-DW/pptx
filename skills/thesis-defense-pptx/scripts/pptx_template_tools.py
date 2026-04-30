from __future__ import annotations

from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


RED = RGBColor(0x83, 0x06, 0x04)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def emu(inches: float) -> int:
    return Inches(inches)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def set_text_preserve_style(shape, text: str, size: float | None = None, bold: bool | None = None, align=None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    old_size = None
    old_bold = None
    old_color = None
    old_name = None
    for old_p in tf.paragraphs:
        for old_r in old_p.runs:
            if old_r.text.strip():
                old_size = old_r.font.size
                old_bold = old_r.font.bold
                old_name = old_r.font.name
                try:
                    old_color = old_r.font.color.rgb
                except Exception:
                    old_color = None
                break
        if old_size is not None or old_color is not None:
            break

    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    if old_name:
        r.font.name = old_name
    if size is not None:
        r.font.size = Pt(size)
    elif old_size is not None:
        r.font.size = old_size
    if bold is not None:
        r.font.bold = bold
    elif old_bold is not None:
        r.font.bold = old_bold
    if old_color is not None:
        r.font.color.rgb = old_color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool = False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align or PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_para(slide, lines: Sequence[str], x: float, y: float, w: float, h: float, size: float, color=BLACK, bold_first: bool = False):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bool(bold_first and i == 0)
        r.font.color.rgb = color
    return box


def rect(slide, x: float, y: float, w: float, h: float, line=RED, fill=None, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, emu(x), emu(y), emu(w), emu(h))
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def tag(slide, text: str, x: float, y: float, w: float, h: float, fill=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    set_text_preserve_style(shape, text, 20, False, PP_ALIGN.CENTER)
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = WHITE
    return shape


def add_pic_fit(slide, src: str | Path, x: float, y: float, w: float, h: float):
    path = Path(src)
    with Image.open(path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio > box_ratio:
        nw = w
        nh = w / img_ratio
    else:
        nh = h
        nw = h * img_ratio
    return slide.shapes.add_picture(str(path), emu(x + (w - nw) / 2), emu(y + (h - nh) / 2), width=emu(nw), height=emu(nh))


def replace_exact_text(slide, mapping: dict[str, str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text in mapping:
                set_text_preserve_style(shape, mapping[text])
