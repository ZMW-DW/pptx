"""Build a tiny PPTX visual template for the minimal example.

The repository does not commit a binary .pptx template. This script generates
one from source so reviewers can see exactly which visual assumptions the
example uses.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "sample_template.pptx"

RED = RGBColor(0x83, 0x06, 0x04)
DARK = RGBColor(0x26, 0x26, 0x26)
LIGHT = RGBColor(0xF7, 0xF3, 0xF1)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, x, y, w, h, size, *, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def nav_bar(slide, active: str) -> None:
    labels = ["背景", "方法", "结果", "总结"]
    for i, label in enumerate(labels):
        x = 0.7 + i * 1.35
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.25), Inches(1.1), Inches(0.38))
        if label == active:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RED
            shape.line.color.rgb = RED
            color = WHITE
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT
            shape.line.color.rgb = RED
            color = RED
        add_text(slide, label, x, 0.27, 1.1, 0.32, 11, bold=True, color=color, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.1)


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    band = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.05))
    band.fill.solid()
    band.fill.fore_color.rgb = RED
    band.line.color.rgb = RED
    add_text(cover, "学校/实验室模板标题区", 0.65, 0.28, 4.5, 0.45, 16, bold=True, color=WHITE)
    add_text(cover, "毕业论文答辩", 0.9, 2.35, 5.4, 0.8, 34, bold=True, color=RED)
    add_text(cover, "题目占位：请替换为论文题目", 0.95, 3.35, 8.5, 0.5, 18)
    add_text(cover, "作者 / 专业 / 指导教师 / 日期", 0.95, 4.15, 7.2, 0.4, 14)

    for title, active in [("研究背景", "背景"), ("方法设计", "方法"), ("实验结果", "结果")]:
        slide = prs.slides.add_slide(blank)
        nav_bar(slide, active)
        add_text(slide, title, 0.75, 0.92, 3.2, 0.45, 24, bold=True, color=RED)
        card(slide, 0.75, 1.55, 5.7, 4.7)
        card(slide, 6.85, 1.55, 5.7, 4.7)
        add_text(slide, "左侧内容占位", 1.05, 1.85, 5.1, 0.45, 18, bold=True, color=RED)
        add_text(slide, "右侧内容占位", 7.15, 1.85, 5.1, 0.45, 18, bold=True, color=RED)

    prs.save(OUT)
    print(OUT)
    return 0


if __name__ == "__main__":
    sys.exit(main())
