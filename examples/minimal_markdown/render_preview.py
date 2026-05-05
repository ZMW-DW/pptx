"""Render a lightweight visual fallback when PowerPoint is unavailable.

python-pptx cannot render PowerPoint exactly like the real Microsoft
PowerPoint engine. Earlier versions of this script tried to "simulate" the
template visually by stacking every shape's text into two fixed cards, which
overlapped badly on dense slides and gave a misleading impression that the
generated PPTX itself was broken.

This rewrite is intentionally conservative:

  * Each slide PNG only shows the slide number, the page title, and a clear
    "preview unavailable" banner so readers do NOT mistake it for a faithful
    render.
  * For real visual QA, run skills/thesis-defense-pptx/scripts/export_pptx_png.ps1
    on Windows + Microsoft PowerPoint. The example's run_example.py prefers
    that path automatically when PowerPoint is available.
"""
from __future__ import annotations

import sys
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

HERE = Path(__file__).resolve().parent
PPTX = HERE / "final.pptx"
OUT_DIR = HERE / "rendered_slides"

W, H = 1280, 720
RED = (0x83, 0x06, 0x04)
DARK = (38, 38, 38)
MUTED = (110, 110, 110)
LIGHT = (247, 243, 241)


def font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf" if bold else "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/PingFang.ttc",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    return ImageFont.load_default()


def first_title(slide) -> str:
    """Return the first non-empty text run on the slide as a rough title."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if text:
            # Take only the first paragraph; further paragraphs are body content.
            return text.splitlines()[0].strip()
    return ""


def render_slide(idx: int, title: str) -> Image.Image:
    img = Image.new("RGB", (W, H), "white")
    draw = ImageDraw.Draw(img)

    # Header band — purely decorative, identifies the example.
    draw.rectangle((0, 0, W, 82), fill=RED)
    draw.text(
        (48, 24),
        "thesis-defense-pptx-skill — fallback preview",
        fill="white",
        font=font(20, bold=True),
    )

    # Slide number + title.
    draw.text((70, 130), f"Slide {idx}", fill=MUTED, font=font(28, bold=True))
    title_text = title or "(no title found)"
    wrapped_title = textwrap.fill(title_text, width=28)
    draw.multiline_text(
        (70, 180), wrapped_title, fill=RED, font=font(38, bold=True), spacing=10
    )

    # Centred banner explaining the fallback so readers don't take this PNG as
    # an accurate render of the slide layout.
    banner_top = 360
    banner_bottom = 560
    draw.rounded_rectangle(
        (70, banner_top, W - 70, banner_bottom), radius=20, fill=LIGHT, outline=RED, width=3
    )
    msg_lines = [
        "Visual preview unavailable in this environment.",
        "python-pptx cannot render PowerPoint exactly.",
        "",
        "For a faithful preview, run on Windows + PowerPoint:",
        "  skills/.../scripts/export_pptx_png.ps1",
    ]
    y = banner_top + 30
    for line in msg_lines:
        is_path = line.lstrip().startswith("skills/")
        f = font(20, bold=False)
        color = DARK if not is_path else RED
        draw.text((100, y), line, fill=color, font=f)
        y += 32

    draw.text((70, 660), f"Slide {idx}", fill=DARK, font=font(18))
    return img


def main() -> int:
    prs = Presentation(str(PPTX))
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for old in OUT_DIR.glob("*.png"):
        old.unlink()
    for idx, slide in enumerate(prs.slides, 1):
        img = render_slide(idx, first_title(slide))
        out = OUT_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        print(out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
