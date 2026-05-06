"""Build a small editable defense deck from thesis.md and sample_template.pptx.

This is intentionally a compact example, not a universal PPT generator. It
shows how an agent can preserve a template's visual language and use the helper
functions shipped in skills/thesis-defense-pptx/scripts/pptx_template_tools.py.
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
REPO_ROOT = HERE.parent.parent
SCRIPT_DIR = REPO_ROOT / "skills" / "thesis-defense-pptx" / "scripts"
sys.path.insert(0, str(SCRIPT_DIR))

from pptx_template_tools import add_para, add_text, rect, tag, write_table, RED, BLACK, WHITE  # noqa: E402

TEMPLATE = HERE / "sample_template.pptx"
THESIS = HERE / "thesis.md"
OUT = HERE / "final.pptx"

MUTED = RGBColor(0x60, 0x66, 0x70)
LIGHT = RGBColor(0xF8, 0xF4, 0xF2)
PALE = RGBColor(0xFB, 0xFB, 0xFB)


def section(text: str, title: str) -> str:
    pattern = re.compile(rf"^##\s+{re.escape(title)}\s*$", re.M)
    match = pattern.search(text)
    if not match:
        return ""
    next_match = re.search(r"^##\s+", text[match.end():], re.M)
    end = match.end() + next_match.start() if next_match else len(text)
    return text[match.end():end].strip()


def compact_lines(text: str, limit: int = 3) -> list[str]:
    lines: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        line = re.sub(r"^\d+[.)、]\s*", "", line)
        line = line.strip("-* \t")
        if not line:
            continue
        if len(line) > 46:
            line = line[:43] + "..."
        lines.append("• " + line)
        if len(lines) >= limit:
            return lines

    sentences = [s.strip() for s in re.split(r"[。.!?]\s*", text) if s.strip()]
    lines = []
    for sentence in sentences[:limit]:
        if len(sentence) > 42:
            sentence = sentence[:39] + "..."
        lines.append("• " + sentence)
    return lines or ["• 内容待补充"]


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def trim_to_n_slides(prs: Presentation, n: int) -> None:
    """Keep only the first ``n`` slide references.

    This does not clear or rewrite slide contents; it only removes references
    from ``sldIdLst`` so PowerPoint renders the first N pages. It is safe for
    complex templates because all native slide shapes remain untouched.
    """
    sld_id_lst = prs.slides._sldIdLst
    for child in list(sld_id_lst)[n:]:
        sld_id_lst.remove(child)


def ensure_slide_count(prs: Presentation, n: int) -> None:
    blank = prs.slide_layouts[6]
    while len(prs.slides) < n:
        prs.slides.add_slide(blank)


def panel(slide, x: float, y: float, w: float, h: float, *, fill=WHITE, line=RED):
    return rect(slide, x, y, w, h, line=line, fill=fill)


def generated_badge(slide, text: str = "Generated from thesis.md") -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.35), Inches(0.28), Inches(2.25), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.color.rgb = RED
    add_text(slide, text, 10.47, 0.31, 2.0, 0.28, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_cover(slide, title: str) -> None:
    clear_slide(slide)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    band.fill.solid()
    band.fill.fore_color.rgb = RED
    band.line.color.rgb = RED
    add_text(slide, "Minimal Markdown 示例", 0.65, 0.28, 4.5, 0.45, 16, bold=True, color=WHITE)
    add_text(slide, "毕业论文答辩", 0.9, 2.3, 5.2, 0.75, 34, bold=True, color=RED)
    add_text(slide, title, 0.95, 3.28, 9.5, 0.65, 20, bold=True, color=BLACK)
    add_text(slide, "可编辑 PPTX / 模板保持 / 质量检查", 0.95, 4.1, 7.6, 0.45, 15)
    add_text(slide, "Author: thesis-defense-pptx-skill", 0.95, 5.25, 5.8, 0.35, 12)


def build_background(slide, lines: list[str]) -> None:
    clear_slide(slide)
    tag(slide, "背景", 0.7, 0.25, 1.1, 0.38)
    add_text(slide, "研究背景", 0.75, 0.92, 3.2, 0.45, 24, bold=True, color=RED)
    rect(slide, 0.75, 1.55, 5.7, 4.7)
    rect(slide, 6.85, 1.55, 5.7, 4.7)
    add_text(slide, "问题场景", 1.05, 1.85, 4.8, 0.35, 18, bold=True, color=RED)
    add_para(slide, lines, 1.05, 2.35, 5.0, 3.25, 15)
    add_text(slide, "模板约束", 7.15, 1.85, 4.8, 0.35, 18, bold=True, color=RED)
    add_para(slide, ["• 不重新设计学校模板", "• 保留封面、导航、配色和卡片样式", "• 生成可编辑 PPTX，而不是图片页"], 7.15, 2.35, 5.0, 3.25, 15)


def build_method(slide, lines: list[str]) -> None:
    clear_slide(slide)
    tag(slide, "方法", 0.7, 0.25, 1.1, 0.38)
    add_text(slide, "方法设计", 0.75, 0.92, 3.2, 0.45, 24, bold=True, color=RED)
    rect(slide, 0.75, 1.55, 11.8, 4.7)
    add_para(slide, lines, 1.1, 1.95, 5.5, 3.6, 15)
    table_shape = slide.shapes.add_table(4, 2, Inches(7.0), Inches(2.0), Inches(4.7), Inches(2.6))
    write_table(
        table_shape.table,
        [
            ["阶段", "产物"],
            ["读取论文", "thesis_context"],
            ["套用模板", "editable pptx"],
            ["质量检查", "contact sheet"],
        ],
        font_size=12,
    )


def build_result(slide, lines: list[str]) -> None:
    clear_slide(slide)
    tag(slide, "结果", 0.7, 0.25, 1.1, 0.38)
    add_text(slide, "结果与总结", 0.75, 0.92, 3.2, 0.45, 24, bold=True, color=RED)
    rect(slide, 0.75, 1.55, 5.7, 4.7)
    rect(slide, 6.85, 1.55, 5.7, 4.7)
    add_text(slide, "最小示例结果", 1.05, 1.85, 4.8, 0.35, 18, bold=True, color=RED)
    add_para(slide, lines, 1.05, 2.35, 5.0, 3.25, 15)
    add_text(slide, "检查项", 7.15, 1.85, 4.8, 0.35, 18, bold=True, color=RED)
    add_para(slide, ["✓ dump shape/text", "✓ scan stale words", "✓ render preview PNG", "✓ contact sheet"], 7.15, 2.35, 5.0, 3.25, 16)


def build_template_cover(slide, title: str) -> None:
    panel(slide, 0.72, 1.05, 11.85, 5.55, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "Minimal Markdown 示例", 1.08, 1.38, 4.1, 0.38, 18, bold=True, color=RED)
    add_text(slide, "毕业论文答辩", 1.08, 2.2, 4.2, 0.55, 30, bold=True, color=BLACK)
    add_text(slide, title, 1.08, 3.05, 10.2, 0.65, 22, bold=True, color=RED)
    add_text(slide, "真实模板保持 / 可编辑 PPTX / 质量检查", 1.08, 4.05, 8.6, 0.38, 15, color=MUTED)
    add_text(slide, "thesis-defense-pptx-skill", 1.08, 5.32, 4.2, 0.35, 13, color=MUTED)


def build_template_agenda(slide) -> None:
    panel(slide, 0.72, 0.88, 11.9, 5.9, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "Contents", 1.05, 1.1, 3.2, 0.52, 28, bold=True, color=RED)
    items = [
        ("01", "研究背景", "为什么需要保留学校模板"),
        ("02", "方法设计", "从论文内容到可编辑 PPTX"),
        ("03", "实验结果", "结构 dump / stale-word scan"),
        ("04", "质量复核", "PowerPoint COM 导出与总览图"),
    ]
    for idx, (num, title, desc) in enumerate(items):
        col = idx % 2
        row = idx // 2
        x = 1.15 + col * 5.4
        y = 2.05 + row * 1.65
        panel(slide, x, y, 4.75, 1.15, fill=LIGHT)
        add_text(slide, num, x + 0.25, y + 0.24, 0.9, 0.48, 24, bold=True, color=RED)
        add_text(slide, title, x + 1.25, y + 0.2, 2.8, 0.35, 17, bold=True, color=BLACK)
        add_text(slide, desc, x + 1.25, y + 0.66, 3.15, 0.28, 11, color=MUTED)


def build_template_section(slide, number: str, title: str, lines: list[str]) -> None:
    panel(slide, 0.72, 0.8, 11.9, 6.0, fill=WHITE)
    generated_badge(slide)
    add_text(slide, f"Part {number}", 1.08, 1.2, 2.1, 0.38, 18, bold=True, color=RED)
    add_text(slide, title, 1.08, 1.9, 6.2, 0.58, 30, bold=True, color=BLACK)
    panel(slide, 1.08, 3.0, 10.1, 2.45, fill=LIGHT)
    add_para(slide, lines, 1.42, 3.36, 9.0, 1.45, 18)


def build_template_background(slide, lines: list[str]) -> None:
    panel(slide, 0.72, 0.8, 11.9, 6.0, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "研究背景", 1.05, 1.05, 3.2, 0.42, 25, bold=True, color=RED)
    panel(slide, 1.05, 1.82, 5.15, 3.95, fill=LIGHT)
    panel(slide, 6.65, 1.82, 5.25, 3.95, fill=PALE)
    add_text(slide, "问题场景", 1.38, 2.1, 2.8, 0.32, 17, bold=True, color=RED)
    add_para(slide, lines, 1.38, 2.58, 4.35, 2.3, 15)
    add_text(slide, "模板保持原则", 6.98, 2.1, 3.4, 0.32, 17, bold=True, color=RED)
    add_para(slide, ["• 不重新设计学校模板", "• 保留封面、配色和导航", "• 生成可编辑 PPTX 而不是图片页"], 6.98, 2.58, 4.35, 2.3, 15)


def build_template_method(slide, lines: list[str]) -> None:
    panel(slide, 0.72, 0.8, 11.9, 6.0, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "方法设计", 1.05, 1.05, 3.2, 0.42, 25, bold=True, color=RED)
    add_para(slide, lines, 1.05, 1.78, 4.95, 3.0, 15)
    flow = [("1", "读取论文"), ("2", "分析模板"), ("3", "生成 PPTX"), ("4", "导出检查")]
    for i, (num, label) in enumerate(flow):
        x = 6.35 + (i % 2) * 2.7
        y = 1.78 + (i // 2) * 1.55
        panel(slide, x, y, 2.25, 1.05, fill=LIGHT)
        add_text(slide, num, x + 0.18, y + 0.22, 0.42, 0.36, 20, bold=True, color=RED)
        add_text(slide, label, x + 0.72, y + 0.3, 1.2, 0.28, 14, bold=True, color=BLACK)
    table_shape = slide.shapes.add_table(4, 2, Inches(6.35), Inches(5.0), Inches(4.95), Inches(1.35))
    write_table(table_shape.table, [["阶段", "产物"], ["读取论文", "thesis_context"], ["套用模板", "editable pptx"], ["质量检查", "contact sheet"]], font_size=10)


def build_template_results(slide, lines: list[str]) -> None:
    panel(slide, 0.72, 0.8, 11.9, 6.0, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "实验结果", 1.05, 1.05, 3.2, 0.42, 25, bold=True, color=RED)
    panel(slide, 1.05, 1.78, 5.15, 4.45, fill=LIGHT)
    add_text(slide, "最小示例结果", 1.38, 2.1, 3.4, 0.32, 17, bold=True, color=RED)
    add_para(slide, lines, 1.38, 2.58, 4.25, 2.4, 15)
    panel(slide, 6.65, 1.78, 4.95, 4.45, fill=PALE)
    add_text(slide, "质量检查输出", 6.98, 2.1, 3.4, 0.32, 17, bold=True, color=RED)
    add_para(slide, ["✓ dump shape/text", "✓ scan stale words", "✓ PowerPoint COM PNG", "✓ README contact sheets"], 6.98, 2.58, 4.0, 2.5, 15)


def build_template_summary(slide) -> None:
    panel(slide, 0.72, 0.8, 11.9, 6.0, fill=WHITE)
    generated_badge(slide)
    add_text(slide, "总结", 1.05, 1.05, 3.2, 0.42, 25, bold=True, color=RED)
    add_text(slide, "模板不是被截图复用，而是作为生成后的可编辑 PPTX 视觉底座。", 1.05, 1.9, 9.9, 0.42, 21, bold=True, color=BLACK)
    add_para(slide, ["• 示例内容来自 thesis.md", "• 输出是 final.pptx，可继续编辑", "• README 图片来自 PowerPoint COM 对 final.pptx 的真实渲染", "• dump / scan / contact sheet 均在生成后文件上执行"], 1.08, 2.85, 9.2, 2.25, 17)


def build_from_real_template(prs: Presentation, title: str, bg_lines: list[str], method_lines: list[str], result_lines: list[str]) -> None:
    ensure_slide_count(prs, 8)
    trim_to_n_slides(prs, 8)
    build_template_cover(prs.slides[0], title)
    build_template_agenda(prs.slides[1])
    build_template_section(prs.slides[2], "01", "研究背景与问题", bg_lines)
    build_template_background(prs.slides[3], bg_lines)
    build_template_section(prs.slides[4], "02", "方法设计", method_lines)
    build_template_method(prs.slides[5], method_lines)
    build_template_results(prs.slides[6], result_lines)
    build_template_summary(prs.slides[7])


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", default=str(TEMPLATE), help="PPTX template to use")
    parser.add_argument("--out", default=str(OUT), help="Output PPTX path")
    parser.add_argument(
        "--real-template",
        action="store_true",
        help="Overlay generated example content on a real institutional template instead of clearing demo slides.",
    )
    args = parser.parse_args()

    template = Path(args.template).expanduser().resolve()
    out = Path(args.out).expanduser().resolve()
    if not template.exists():
        raise SystemExit(f"template not found: {template}")
    text = THESIS.read_text(encoding="utf-8")
    prs = Presentation(str(template))

    title = text.splitlines()[0].lstrip("# ").strip()
    bg_lines = compact_lines(section(text, "研究背景"), 3)
    method_lines = compact_lines(section(text, "方法设计"), 4)
    result_text = section(text, "实验结果") + "\n" + section(text, "结论")
    result_lines = compact_lines(result_text, 4)

    if args.real_template:
        build_from_real_template(prs, title, bg_lines, method_lines, result_lines)
    else:
        while len(prs.slides) < 4:
            prs.slides.add_slide(prs.slide_layouts[6])
        build_cover(prs.slides[0], title)
        build_background(prs.slides[1], bg_lines)
        build_method(prs.slides[2], method_lines)
        build_result(prs.slides[3], result_lines)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
